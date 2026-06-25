import fitz  # PyMuPDF
import pytesseract
import io  # noqa: F401  (kept for potential future callers)
import os  # noqa: F401
import re
import logging
import httpx
from PIL import Image
from dataclasses import dataclass, field
from typing import List, Tuple, Optional
from pptx import Presentation
from docx2python import docx2python
from app.config import settings
from app.models.document import Document

logger = logging.getLogger(__name__)

# Point pytesseract at the Windows installation path configured in .env
pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd

# Patterns that look like LaTeX / AsciiMath equations.
# Ordered from most-specific (display) to least-specific (inline).
EQUATION_PATTERNS = [
    r'\$\$[\s\S]+?\$\$',                            # $$...$$  display math
    r'\$[^\$\n]+?\$',                               # $...$    inline math
    r'\\\[[\s\S]+?\\\]',                            # \[...\]  display math
    r'\\\([\s\S]+?\\\)',                             # \(...\)  inline math
    r'\\begin\{equation\*?\}[\s\S]+?\\end\{equation\*?\}',
    r'\\begin\{align\*?\}[\s\S]+?\\end\{align\*?\}',
]


# ── Data classes ──────────────────────────────────────────────────────────────

@dataclass
class TextBlock:
    text: str
    page_number: int
    block_type: str = "text"
    bbox: Optional[Tuple] = None


@dataclass
class ImageBlock:
    image_bytes: bytes
    page_number: int
    caption: str = ""
    block_type: str = "figure"


@dataclass
class ParsedDocument:
    document_id: str
    file_type: str
    text_blocks: List[TextBlock] = field(default_factory=list)
    image_blocks: List[ImageBlock] = field(default_factory=list)
    page_count: int = 0


# ── Parser ────────────────────────────────────────────────────────────────────

class DocumentParser:
    """Parses PDF / PPTX / DOCX / image files into TextBlock and ImageBlock lists."""

    def parse_document(self, doc: Document) -> ParsedDocument:
        parsers = {
            "pdf": self._parse_pdf,
            "pptx": self._parse_pptx,
            "docx": self._parse_docx,
            "image": self._parse_image,
        }
        parser = parsers.get(doc.file_type)
        if not parser:
            raise ValueError(f"Unknown file type: {doc.file_type}")
        return parser(doc)

    # ── Equation extraction ───────────────────────────────────────────────────

    def _extract_equations(self, text: str) -> List[str]:
        """Return unique equation strings found in text."""
        found: List[str] = []
        for pattern in EQUATION_PATTERNS:
            found.extend(re.findall(pattern, text, re.DOTALL))
        return list(set(found))

    # ── PDF ───────────────────────────────────────────────────────────────────

    def _parse_pdf(self, doc: Document) -> ParsedDocument:
        parsed = ParsedDocument(document_id=str(doc.id), file_type="pdf")
        try:
            pdf = fitz.open(doc.file_path)
            parsed.page_count = len(pdf)
            seen_xrefs: set = set()  # prevent duplicate images across pages
            for page_num, page in enumerate(pdf, 1):
                try:
                    # ── Text: one TextBlock per PyMuPDF paragraph block ──────
                    # Using "blocks" gives paragraph-level granularity, which
                    # produces far more (and shorter) chunks than whole-page text.
                    for block in page.get_text("blocks"):
                        block_type = block[6]   # 0 = text, 1 = image
                        if block_type != 0:
                            continue
                        block_text = block[4].strip()
                        if len(block_text) >= 5:
                            parsed.text_blocks.append(
                                TextBlock(text=block_text, page_number=page_num)
                            )

                    # ── Images: primary path via xref table ──────────────────
                    page_xrefs: list = []
                    for img_info in page.get_images(full=True):
                        xref = img_info[0]
                        if xref in seen_xrefs:
                            continue
                        seen_xrefs.add(xref)
                        page_xrefs.append(xref)
                        try:
                            base_image = pdf.extract_image(xref)
                            w = base_image["width"]
                            h = base_image["height"]
                            if w >= 60 and h >= 60:
                                parsed.image_blocks.append(
                                    ImageBlock(
                                        image_bytes=base_image["image"],
                                        page_number=page_num,
                                    )
                                )
                        except Exception as e:
                            logger.debug("Failed to extract xref %d on page %d: %s", xref, page_num, e)

                    # ── Images: fallback via dict blocks (inline / form XObjects
                    #    that don't appear in get_images) ──────────────────────
                    if not page_xrefs:
                        for block in page.get_text("dict")["blocks"]:
                            if block.get("type") != 1:
                                continue
                            img_bytes = block.get("image", b"")
                            w = block.get("width", 0)
                            h = block.get("height", 0)
                            if w >= 60 and h >= 60 and len(img_bytes) > 1000:
                                parsed.image_blocks.append(
                                    ImageBlock(
                                        image_bytes=img_bytes,
                                        page_number=page_num,
                                    )
                                )

                except Exception as e:
                    logger.warning("Failed to parse page %d: %s", page_num, e)
            pdf.close()
        except Exception as e:
            logger.error("PDF parse failed: %s", e)
            raise
        return parsed

    # ── PPTX ──────────────────────────────────────────────────────────────────

    def _parse_pptx(self, doc: Document) -> ParsedDocument:
        parsed = ParsedDocument(document_id=str(doc.id), file_type="pptx")
        try:
            prs = Presentation(doc.file_path)
            parsed.page_count = len(prs.slides)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_text.append(t)
                    if hasattr(shape, "image"):
                        try:
                            img_bytes = shape.image.blob
                            if len(img_bytes) > 5000:
                                parsed.image_blocks.append(
                                    ImageBlock(
                                        image_bytes=img_bytes,
                                        page_number=slide_num,
                                    )
                                )
                        except Exception:
                            pass
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_text.append(f"[Notes] {notes}")
                if slide_text:
                    parsed.text_blocks.append(
                        TextBlock(text="\n".join(slide_text), page_number=slide_num)
                    )
        except Exception as e:
            logger.error("PPTX parse failed: %s", e)
            raise
        return parsed

    # ── DOCX ──────────────────────────────────────────────────────────────────

    def _parse_docx(self, doc: Document) -> ParsedDocument:
        parsed = ParsedDocument(
            document_id=str(doc.id), file_type="docx", page_count=1
        )
        try:
            result = docx2python(doc.file_path)
            all_text: List[str] = []
            for section in result.body:
                for para in section:
                    for run in para:
                        t = run.strip() if isinstance(run, str) else ""
                        if t:
                            all_text.append(t)
            if all_text:
                parsed.text_blocks.append(
                    TextBlock(text="\n".join(all_text), page_number=1)
                )
        except Exception as e:
            logger.error("DOCX parse failed: %s", e)
            raise
        return parsed

    # ── Image ─────────────────────────────────────────────────────────────────

    def _parse_image(self, doc: Document) -> ParsedDocument:
        parsed = ParsedDocument(
            document_id=str(doc.id), file_type="image", page_count=1
        )
        try:
            image = Image.open(doc.file_path)
            text = pytesseract.image_to_string(image)
            if text.strip():
                # Optionally augment with MathPix equation OCR
                if settings.mathpix_app_id and settings.mathpix_app_key:
                    latex = self._mathpix_ocr(doc.file_path)
                    if latex:
                        text += "\n" + latex
                parsed.text_blocks.append(
                    TextBlock(text=text, page_number=1)
                )
            with open(doc.file_path, "rb") as f:
                img_bytes = f.read()
            parsed.image_blocks.append(
                ImageBlock(image_bytes=img_bytes, page_number=1)
            )
        except Exception as e:
            logger.error("Image parse failed: %s", e)
            raise
        return parsed

    # ── MathPix helper ────────────────────────────────────────────────────────

    def _mathpix_ocr(self, file_path: str) -> str:
        try:
            import base64
            with open(file_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode()
            response = httpx.post(
                "https://api.mathpix.com/v3/text",
                headers={
                    "app_id": settings.mathpix_app_id,
                    "app_key": settings.mathpix_app_key,
                },
                json={
                    "src": f"data:image/jpeg;base64,{img_b64}",
                    "formats": ["latex_simplified"],
                },
                timeout=30,
            )
            return response.json().get("latex_simplified", "")
        except Exception as e:
            logger.warning("MathPix OCR failed: %s", e)
            return ""
